"""
PPTX 文档解析器 — 单一职责：将 PPT 解析为增强文本（文本 + HTML 表格 + 图片描述/URL）。

使用 python-pptx 提取幻灯片文本、表格、图表和内嵌图片：
    - 文本：shape.text_frame（文本框、占位符）
    - 表格：shape.has_table → HTML <table>
    - 图表：shape.has_chart → 提取数据点为文本
    - 图片：shape.shape_type == PICTURE → 上传 MinIO / VLM 描述
    - GROUP：递归遍历组合形状的子形状（文本/表格/图表/图片）

每个 slide 输出为 <h2>幻灯片 N: 标题</h2> + 内容，
chunker 的 _split_html 天然按 slide 分块。

P0-P2 增强（对齐 PDF/DOCX/XLSX 解析器）：
    - P0: GROUP 组合形状递归提取表格/图表/图片（修复数据丢失）
    - P0: 配置读取用 _bool()/_int() 安全转换（修复 MagicMock 测试）
    - P1: 图片上传 MinIO + 小图过滤（复用 image_storage.py）
    - P1: 分页分隔符支持（page_separator）
    - P1: 表格列宽对齐（合并单元格补空）
    - P2: 图表数据提取（has_chart → 数据点文本）
    - P2: 清理 picture_val 冗余逻辑 + 日志统计修正

遵循优雅降级：
    - python-pptx 未安装 → 返回空字符串；
    - VLM 不可用 → 跳过图片描述；
    - MinIO 不可用 → 跳过图片上传，降级为 VLM 描述；
    - 单个形状解析异常 → 跳过该形状，继续处理其余。
"""

from __future__ import annotations

import asyncio
from typing import Any

from app.config import get_settings
from app.document.base import DocumentParser, ParsedSection
from app.utils.logger import get_logger

log = get_logger(__name__)

# VLM 并发控制
_VLM_SEMAPHORE_LIMIT: int = 3
_IMAGE_PROMPT: str = (
    "请用一句话描述这张图片的内容，"
    "重点关注图表、数据、文字和关键信息，便于后续检索。"
)

# python-pptx 形状类型常量
_SHAPE_TYPE_GROUP = 6
_SHAPE_TYPE_PICTURE = 13


class PPTXParser(DocumentParser):
    """PPTX 解析器 — python-pptx 文本 + 表格 + 图表 + 图片 URL/VLM + GROUP 递归。"""

    async def parse(self, file_path: str) -> str:
        """解析 PPTX 文档，返回增强文本。

        Args:
            file_path: PPTX 文件路径。

        Returns:
            增强文本（<h2>幻灯片标题</h2> + 文本 + HTML 表格 + 图片描述/URL）。
            python-pptx 未安装或解析失败时返回空字符串。
        """
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            log.warning("pptx.parser_skipped", reason="python-pptx not installed")
            return ""

        try:
            prs = Presentation(file_path)
        except Exception as exc:
            log.warning("pptx.open_failed", file_path=file_path, error=str(exc))
            return ""

        settings = get_settings()
        # P0: 使用 _bool()/_int() 安全转换（处理 MagicMock 测试场景）
        image_enabled = self._bool(
            getattr(settings, "PPTX_IMAGE_EXTRACTION_ENABLED", True), True
        )
        max_images = self._int(getattr(settings, "PPTX_IMAGE_MAX_PER_DOC", 50), 50)
        # P1: 图片上传 + 小图过滤配置
        image_upload_enabled = self._bool(
            getattr(settings, "PPTX_IMAGE_UPLOAD_ENABLED", False), False
        )
        image_min_size = self._int(getattr(settings, "PPTX_IMAGE_MIN_SIZE", 50), 50)

        sections: list[ParsedSection] = []
        image_count = 0
        semaphore = asyncio.Semaphore(_VLM_SEMAPHORE_LIMIT)

        for slide_num, slide in enumerate(prs.slides):
            slide_sections = self._extract_slide_text(slide, slide_num)
            sections.extend(slide_sections)

            # P0: 递归提取表格（含 GROUP 组合形状内的表格）
            table_sections = self._extract_tables_recursive(slide, slide_num)
            sections.extend(table_sections)

            # P0: 递归提取图表（含 GROUP 组合形状内的图表）
            chart_sections = self._extract_charts_recursive(slide, slide_num)
            sections.extend(chart_sections)

            # P1: 递归提取图片（含 GROUP + 图片上传 + 小图过滤）
            if image_enabled and image_count < max_images:
                img_sections, img_count = await self._extract_images_recursive(
                    slide,
                    slide_num,
                    MSO_SHAPE_TYPE,
                    semaphore,
                    max_images - image_count,
                    image_upload_enabled=image_upload_enabled,
                    image_min_size=image_min_size,
                )
                sections.extend(img_sections)
                image_count += img_count

            # 提取演讲者备注
            note_section = self._extract_notes(slide, slide_num)
            if note_section:
                sections.append(note_section)

        log.info(
            "pptx.parsed",
            file_path=file_path,
            slides=len(prs.slides),
            sections=len(sections),
            tables=sum(1 for s in sections if s.kind == "table"),
            charts=sum(1 for s in sections if s.kind == "text" and "[图表数据]" in s.content),
            images=sum(1 for s in sections if s.kind in ("image_desc", "image_url")),
        )

        # P1: 分页分隔符支持
        page_sep_raw = getattr(settings, "PAGE_SEPARATOR", "")
        page_separator = page_sep_raw if isinstance(page_sep_raw, str) else ""
        return self.sections_to_text(sections, page_separator=page_separator)

    def _extract_slide_text(self, slide: Any, slide_num: int) -> list[ParsedSection]:
        """提取幻灯片文本，包装为 <h2> 标题块。

        尝试从 slide 的第一个标题占位符提取标题，找不到则用"幻灯片 N"。
        递归遍历 GROUP 组合形状的子形状，提取组合内的文本框。
        """
        title = self._get_slide_title(slide) or f"幻灯片 {slide_num + 1}"
        title_escaped = self._escape_html(title)
        text_parts: list[str] = [f"<h2>{title_escaped}</h2>"]

        # 递归遍历所有形状（包括 GROUP 组合形状的子形状）
        for shape in slide.shapes:
            self._collect_shape_text(shape, text_parts)

        content = "\n".join(text_parts).strip()
        if not content:
            return []

        return [ParsedSection(kind="text", content=content, page=slide_num)]

    def _collect_shape_text(
        self, shape: Any, text_parts: list[str], depth: int = 0
    ) -> None:
        """递归收集形状文本 — 处理 GROUP 组合形状。

        GROUP 形状本身没有 text_frame，但其子形状可能有。
        递归遍历子形状提取文本，同时跳过表格、图片和图表（由专用方法处理）。

        Args:
            shape: pptx Shape 对象。
            text_parts: 文本收集列表（可变引用）。
            depth: 递归深度（防止无限循环，上限 5 层）。
        """
        if depth > 5:
            return

        # 跳过表格（由 _extract_tables_recursive 处理）
        if shape.has_table:
            return

        # 跳过图表（由 _extract_charts_recursive 处理）
        if hasattr(shape, "has_chart") and shape.has_chart:
            return

        # 检查形状类型
        shape_type_val = self._get_shape_type_value(shape)

        # 跳过图片（PICTURE = 13，由 _extract_images_recursive 处理）
        if shape_type_val == _SHAPE_TYPE_PICTURE:
            return

        # GROUP 形状 — 递归遍历子形状
        if shape_type_val == _SHAPE_TYPE_GROUP:
            try:
                for child_shape in shape.shapes:
                    self._collect_shape_text(child_shape, text_parts, depth + 1)
            except Exception as exc:
                log.debug("pptx.group_traverse_failed", depth=depth, error=str(exc))
            return

        # 普通形状 — 提取文本
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                text_parts.append(self._escape_html(text))

    @staticmethod
    def _get_shape_type_value(shape: Any) -> int | None:
        """安全获取形状类型值。

        Args:
            shape: pptx Shape 对象。

        Returns:
            形状类型整数值（如 6=GROUP, 13=PICTURE），获取失败返回 None。
        """
        try:
            shape_type = shape.shape_type
            if shape_type is not None and hasattr(shape_type, "value"):
                return shape_type.value
        except Exception:
            pass
        return None

    @staticmethod
    def _get_slide_title(slide: Any) -> str:
        """尝试从幻灯片提取标题文本。"""
        try:
            if slide.shapes.title:
                return slide.shapes.title.text.strip()
        except Exception:
            pass

        # 尝试从 placeholders 提取
        try:
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.idx == 0:  # 标题占位符
                    return placeholder.text.strip()
        except Exception:
            pass

        return ""

    @staticmethod
    def _extract_notes(slide: Any, slide_num: int) -> ParsedSection | None:
        """提取幻灯片的演讲者备注。

        备注存储在 slide.notes_slide.notes_text_frame 中，
        包含演讲者补充说明，对 RAG 检索有较高价值。

        Args:
            slide: pptx Slide 对象。
            slide_num: 幻灯片编号。

        Returns:
            备注的 ParsedSection，无备注时返回 None。
        """
        try:
            # has_notes_slide 在某些版本中不可用，用 try/except 保护
            if not slide.has_notes_slide:
                return None
        except Exception:
            pass

        try:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text.strip()
            if not notes_text:
                return None

            return ParsedSection(
                kind="text",
                content=f"[演讲者备注]\n{notes_text}",
                page=slide_num,
            )
        except Exception as exc:
            log.debug("pptx.notes_extract_failed", slide=slide_num, error=str(exc))
            return None

    # ============================================================
    # P0: 表格递归提取（含 GROUP 组合形状）
    # ============================================================

    def _extract_tables_recursive(
        self, slide: Any, slide_num: int
    ) -> list[ParsedSection]:
        """递归提取幻灯片中的表格 — 含 GROUP 组合形状内的表格。

        P0 修复：原 _extract_tables 只遍历顶层 shapes，
        GROUP 组合形状内的表格完全丢失。本方法递归遍历所有形状，
        包括 GROUP 的子形状，提取所有表格。

        Args:
            slide: pptx Slide 对象。
            slide_num: 幻灯片编号。

        Returns:
            表格 ParsedSection 列表。
        """
        sections: list[ParsedSection] = []
        for shape in slide.shapes:
            self._collect_tables(shape, slide_num, sections, depth=0)
        return sections

    def _collect_tables(
        self,
        shape: Any,
        slide_num: int,
        sections: list[ParsedSection],
        depth: int = 0,
    ) -> None:
        """递归收集表格 — 处理 GROUP 组合形状。

        Args:
            shape: pptx Shape 对象。
            slide_num: 幻灯片编号。
            sections: 表格 ParsedSection 收集列表。
            depth: 递归深度（上限 5 层）。
        """
        if depth > 5:
            return

        # 检测表格
        if shape.has_table:
            try:
                table = shape.table
                rows: list[list[str | None]] = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)

                html = self._rows_to_html(rows)
                if html:
                    sections.append(
                        ParsedSection(kind="table", content=html, page=slide_num)
                    )
            except Exception as exc:
                log.debug("pptx.table_extract_failed", slide=slide_num, error=str(exc))
            return

        # GROUP 形状 — 递归遍历子形状
        shape_type_val = self._get_shape_type_value(shape)
        if shape_type_val == _SHAPE_TYPE_GROUP:
            try:
                for child_shape in shape.shapes:
                    self._collect_tables(child_shape, slide_num, sections, depth + 1)
            except Exception as exc:
                log.debug("pptx.group_table_traverse_failed", depth=depth, error=str(exc))

    # ============================================================
    # P2: 图表数据提取（含 GROUP 组合形状）
    # ============================================================

    def _extract_charts_recursive(
        self, slide: Any, slide_num: int
    ) -> list[ParsedSection]:
        """递归提取幻灯片中的图表数据 — 含 GROUP 组合形状内的图表。

        P2 增强：原解析器完全忽略图表（has_chart），
        导致 PPT 中的柱状图/饼图/折线图数据全部丢失。
        本方法提取图表的数据点为文本，便于 RAG 检索。

        Args:
            slide: pptx Slide 对象。
            slide_num: 幻灯片编号。

        Returns:
            图表数据 ParsedSection 列表。
        """
        sections: list[ParsedSection] = []
        for shape in slide.shapes:
            self._collect_charts(shape, slide_num, sections, depth=0)
        return sections

    def _collect_charts(
        self,
        shape: Any,
        slide_num: int,
        sections: list[ParsedSection],
        depth: int = 0,
    ) -> None:
        """递归收集图表数据 — 处理 GROUP 组合形状。

        Args:
            shape: pptx Shape 对象。
            slide_num: 幻灯片编号。
            sections: 图表数据 ParsedSection 收集列表。
            depth: 递归深度（上限 5 层）。
        """
        if depth > 5:
            return

        # 检测图表
        if hasattr(shape, "has_chart") and shape.has_chart:
            try:
                chart_text = self._extract_chart_data(shape.chart)
                if chart_text:
                    sections.append(
                        ParsedSection(
                            kind="text",
                            content=f"[图表数据]\n{chart_text}",
                            page=slide_num,
                        )
                    )
            except Exception as exc:
                log.debug("pptx.chart_extract_failed", slide=slide_num, error=str(exc))
            return

        # GROUP 形状 — 递归遍历子形状
        shape_type_val = self._get_shape_type_value(shape)
        if shape_type_val == _SHAPE_TYPE_GROUP:
            try:
                for child_shape in shape.shapes:
                    self._collect_charts(child_shape, slide_num, sections, depth + 1)
            except Exception as exc:
                log.debug("pptx.group_chart_traverse_failed", depth=depth, error=str(exc))

    @staticmethod
    def _extract_chart_data(chart: Any) -> str:
        """提取图表数据点为文本。

        支持 plot.series.categories 和 series.values，
        将数据点格式化为"类别: 值"的文本列表。

        Args:
            chart: pptx Chart 对象。

        Returns:
            图表数据文本。提取失败返回空字符串。
        """
        try:
            lines: list[str] = []
            plot = chart.plots[0]

            # 获取类别标签
            categories = list(plot.categories) if plot.categories else []

            for series in plot.series:
                series_name = series.name or "系列"
                values = list(series.values) if series.values else []

                lines.append(f"系列: {series_name}")
                for idx, val in enumerate(values):
                    cat = categories[idx] if idx < len(categories) else f"项{idx + 1}"
                    lines.append(f"  {cat}: {val}")
                lines.append("")  # 空行分隔

            return "\n".join(lines).strip()
        except Exception as exc:
            log.debug("pptx.chart_data_failed", error=str(exc))
            return ""

    # ============================================================
    # P1: 图片递归提取（含 GROUP + 上传 MinIO + 小图过滤）
    # ============================================================

    async def _extract_images_recursive(
        self,
        slide: Any,
        slide_num: int,
        mso_shape_type: Any,
        semaphore: asyncio.Semaphore,
        remaining: int,
        image_upload_enabled: bool = False,
        image_min_size: int = 50,
    ) -> tuple[list[ParsedSection], int]:
        """递归提取幻灯片内嵌图片 — 含 GROUP + 上传 MinIO + 小图过滤。

        P0 修复：递归遍历 GROUP 组合形状，提取所有图片。
        P1 增强：复用 image_storage.py 实现图片上传和小图过滤。

        Args:
            slide: pptx Slide 对象。
            slide_num: 幻灯片编号。
            mso_shape_type: MSO_SHAPE_TYPE 枚举类。
            semaphore: VLM 并发信号量。
            remaining: 本文档剩余可提取图片数。
            image_upload_enabled: 是否上传图片到 MinIO。
            image_min_size: 最小尺寸阈值（宽或高小于此值跳过）。

        Returns:
            (图片 ParsedSection 列表, 实际提取的图片数)
        """
        # P2: 清理冗余逻辑 — 直接用常量
        picture_val = _SHAPE_TYPE_PICTURE

        # 递归收集所有图片形状（含 GROUP 内的）
        image_shapes: list[Any] = []
        for shape in slide.shapes:
            self._collect_image_shapes(shape, image_shapes, remaining - len(image_shapes), depth=0)

        if not image_shapes:
            return [], 0

        image_shapes = image_shapes[:remaining]
        sections: list[ParsedSection] = []
        count = 0

        async def process_image(idx: int, shape: Any) -> ParsedSection | None:
            """处理单张图片 — 提取、过滤、上传/描述。"""
            try:
                img_bytes = shape.image.blob
                ext = "png"
                try:
                    ext = shape.image.ext
                except Exception:
                    pass
            except Exception as exc:
                log.debug("pptx.image_extract_failed", slide=slide_num, error=str(exc))
                return None

            if not img_bytes:
                return None

            # 标准化扩展名
            std_ext = ext.lower().lstrip(".")
            if std_ext == "jpg":
                std_ext = "jpeg"
            mime_type = f"image/{std_ext}"

            # P1: 小图过滤（仅在尺寸已知时过滤）
            img_width = 0
            img_height = 0
            if image_min_size > 0:
                from app.document.image_storage import get_image_dimensions

                img_width, img_height = get_image_dimensions(img_bytes, std_ext)
                if img_width > 0 and img_height > 0:
                    if img_width < image_min_size or img_height < image_min_size:
                        log.debug(
                            "pptx.images.filtered_small",
                            width=img_width,
                            height=img_height,
                            min_size=image_min_size,
                            slide=slide_num,
                        )
                        return None

            # VLM 描述
            vlm_desc = ""
            async with semaphore:
                vlm_desc = await self._vlm_describe(img_bytes, mime_type)

            desc_text = ""
            if vlm_desc and vlm_desc.strip():
                desc_text = f"[图片描述: {vlm_desc.strip()}]"

            # P1: 图片上传模式
            if image_upload_enabled:
                from app.document.image_storage import upload_image

                url = await upload_image(
                    image_bytes=img_bytes,
                    ext=std_ext,
                    doc_id="pptx",
                    page=slide_num,
                    idx=idx,
                    min_size=0,  # 已过滤
                    width=img_width,
                    height=img_height,
                )
                if url:
                    return ParsedSection(
                        kind="image_url",
                        content=desc_text,
                        page=slide_num,
                        image_url=url,
                    )
                # 上传失败，降级为描述模式
                if desc_text:
                    return ParsedSection(
                        kind="image_desc",
                        content=desc_text,
                        page=slide_num,
                    )
                return None

            # 仅描述模式
            if not desc_text:
                return None
            return ParsedSection(
                kind="image_desc",
                content=desc_text,
                page=slide_num,
            )

        tasks = [process_image(idx, shape) for idx, shape in enumerate(image_shapes)]
        results = await asyncio.gather(*tasks, return_exceptions=True)

        for result in results:
            if isinstance(result, ParsedSection):
                sections.append(result)
                count += 1
            elif isinstance(result, Exception):
                log.warning("pptx.image_process_error", error=str(result))

        return sections, count

    def _collect_image_shapes(
        self,
        shape: Any,
        image_shapes: list[Any],
        remaining: int,
        depth: int = 0,
    ) -> None:
        """递归收集图片形状 — 处理 GROUP 组合形状。

        Args:
            shape: pptx Shape 对象。
            image_shapes: 图片形状收集列表。
            remaining: 剩余可收集数量。
            depth: 递归深度（上限 5 层）。
        """
        if depth > 5 or remaining <= 0:
            return

        shape_type_val = self._get_shape_type_value(shape)

        # 图片形状
        if shape_type_val == _SHAPE_TYPE_PICTURE:
            image_shapes.append(shape)
            return

        # GROUP 形状 — 递归遍历子形状
        if shape_type_val == _SHAPE_TYPE_GROUP:
            try:
                for child_shape in shape.shapes:
                    if len(image_shapes) >= remaining:
                        break
                    self._collect_image_shapes(
                        child_shape, image_shapes, remaining - len(image_shapes), depth + 1
                    )
            except Exception as exc:
                log.debug("pptx.group_image_traverse_failed", depth=depth, error=str(exc))

    # ============================================================
    # 工具方法
    # ============================================================

    @staticmethod
    def _rows_to_html(rows: list[list[str | None]]) -> str:
        """将二维数据转为 HTML <table> 标签 — P1: 列宽对齐。

        第一行视为表头（<th>），其余为数据行（<td>）。
        P1 增强：自动补齐不等长行为最大列数（合并单元格场景）。
        """
        if not rows:
            return ""

        # P1: 列宽对齐 — 找出最大列数并补齐
        max_cols = max(len(r) for r in rows) if rows else 0
        padded_rows = [list(r) + [""] * (max_cols - len(r)) for r in rows]

        lines: list[str] = ["<table>"]

        for i, row in enumerate(padded_rows):
            lines.append("<tr>")
            tag = "th" if i == 0 else "td"
            for cell in row:
                cell_text = (cell or "").strip()
                cell_text = PPTXParser._escape_html(cell_text)
                lines.append(f"<{tag}>{cell_text}</{tag}>")
            lines.append("</tr>")

        lines.append("</table>")
        return "\n".join(lines)

    @staticmethod
    def _escape_html(text: str) -> str:
        """转义 HTML 特殊字符。"""
        if not text:
            return ""
        return (
            text.replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
        )

    async def _vlm_describe(self, image: bytes, mime_type: str) -> str:
        """调用 VLM 理解图片内容 — 延迟导入，不可用时返回空字符串。"""
        try:
            from app.vlm.provider import get_vision_provider

            vlm = get_vision_provider()
            return await vlm.understand(
                image=image,
                prompt=_IMAGE_PROMPT,
                mime_type=mime_type,
            )
        except ImportError:
            log.warning("pptx.vlm_not_available")
            return ""
        except Exception as exc:
            log.warning("pptx.vlm_error", error=str(exc))
            return ""
